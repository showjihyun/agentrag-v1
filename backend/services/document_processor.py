"""
Document Processing Service for extracting and chunking text from various file formats.

Supports PDF, TXT, DOCX, HWP, and HWPX file formats with configurable chunking
and comprehensive error handling.
"""

import logging
import uuid
from typing import List, Dict, Any, Optional, BinaryIO
from datetime import datetime
import io
import zipfile
import xml.etree.ElementTree as ET

# File processing libraries
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
import olefile

# Internal models
from backend.models.document import Document, TextChunk

logger = logging.getLogger(__name__)


class DocumentProcessingError(Exception):
    """Custom exception for document processing errors."""

    pass


class DocumentProcessor:
    """
    Service for processing documents: extraction, chunking, and metadata management.

    Features:
    - Multi-format support (PDF, TXT, DOCX, HWP, HWPX)
    - Configurable text chunking with overlap
    - File size validation
    - Metadata extraction
    - Comprehensive error handling
    """

    # Supported file types
    SUPPORTED_TYPES = {
        "pdf": "application/pdf",
        "txt": "text/plain",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "hwp": "application/x-hwp",
        "hwpx": "application/x-hwpx",
        "md": "text/markdown",
        "ppt": "application/vnd.ms-powerpoint",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "xls": "application/vnd.ms-excel",
        "csv": "text/csv",
        "json": "application/json",
        # Image formats (OCR support)
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "bmp": "image/bmp",
        "webp": "image/webp",
    }

    def __init__(
        self,
        chunk_size: int = None,
        chunk_overlap: int = None,
        max_file_size: int = None,
        chunking_strategy: str = None,
    ):
        """
        Initialize DocumentProcessor.

        Args:
            chunk_size: Target size for text chunks in characters (default: from config)
            chunk_overlap: Number of overlapping characters between chunks (default: from config)
            max_file_size: Maximum allowed file size in bytes (default: from config)
            chunking_strategy: Strategy for text chunking (default: from config)
                - "semantic": Semantic similarity-based (best quality)
                - "sentence": Sentence boundary-based (good balance)
                - "paragraph": Paragraph boundary-based (fast)
                - "heading": Heading-based (for structured docs)
                - "fixed": Fixed-size with boundaries (fallback)

        Raises:
            ValueError: If parameters are invalid
        """
        # Load from config if not provided
        from backend.config import settings
        
        # Ensure values are integers (config might return strings from env)
        self.chunk_size = int(chunk_size if chunk_size is not None else settings.CHUNK_SIZE)
        self.chunk_overlap = int(chunk_overlap if chunk_overlap is not None else settings.CHUNK_OVERLAP)
        self.max_file_size = int(max_file_size if max_file_size is not None else settings.MAX_FILE_SIZE)
        self.chunking_strategy = str(chunking_strategy if chunking_strategy is not None else settings.CHUNKING_STRATEGY)
        
        # Validate parameters
        if self.chunk_size <= 0:
            raise ValueError("chunk_size must be positive")
        if self.chunk_overlap < 0:
            raise ValueError("chunk_overlap cannot be negative")
        if self.chunk_overlap >= self.chunk_size:
            raise ValueError("chunk_overlap must be less than chunk_size")
        if self.max_file_size <= 0:
            raise ValueError("max_file_size must be positive")

        # Initialize semantic chunker
        from backend.services.semantic_chunker import get_semantic_chunker

        self.semantic_chunker = get_semantic_chunker(
            strategy=self.chunking_strategy,  # Use self.chunking_strategy, not chunking_strategy
            target_size=self.chunk_size,      # Use self.chunk_size, not chunk_size
            overlap=self.chunk_overlap        # Use self.chunk_overlap, not chunk_overlap
        )

        # Initialize enhanced processor with contextual retrieval
        from backend.services.enhanced_document_processor import (
            get_enhanced_document_processor,
        )

        self.enhanced_processor = get_enhanced_document_processor()

        # Initialize metadata extractor
        from backend.services.metadata_extractor import get_metadata_extractor

        self.metadata_extractor = get_metadata_extractor()

        logger.info(
            f"DocumentProcessor initialized with Contextual Retrieval: chunk_size={chunk_size}, "
            f"overlap={chunk_overlap}, max_size={max_file_size}, "
            f"strategy={chunking_strategy}"
        )

    def detect_file_type(self, filename: str) -> str:
        """
        Detect file type from filename extension.

        Args:
            filename: Name of the file

        Returns:
            str: File type (pdf, txt, docx, hwp, hwpx)

        Raises:
            ValueError: If file type is not supported
        """
        if not filename:
            raise ValueError("filename cannot be empty")

        # Extract extension
        extension = filename.lower().split(".")[-1] if "." in filename else ""

        if extension not in self.SUPPORTED_TYPES:
            raise ValueError(
                f"Unsupported file type: .{extension}. "
                f"Supported types: {', '.join(self.SUPPORTED_TYPES.keys())}"
            )

        logger.debug(f"Detected file type: {extension} for {filename}")
        return extension

    def validate_file_size(self, file_size: int) -> None:
        """
        Validate that file size is within limits.

        Args:
            file_size: Size of file in bytes

        Raises:
            ValueError: If file size exceeds maximum
        """
        if file_size <= 0:
            raise ValueError("file_size must be positive")

        if file_size > self.max_file_size:
            raise ValueError(
                f"File size {file_size} bytes exceeds maximum "
                f"{self.max_file_size} bytes ({self.max_file_size / 1024 / 1024:.1f}MB)"
            )

        logger.debug(f"File size validated: {file_size} bytes")

    def extract_text_from_pdf(self, file_content: bytes) -> str:
        """
        Extract text from PDF file using opendataloader-pdf.

        Args:
            file_content: PDF file content as bytes

        Returns:
            str: Extracted text

        Raises:
            DocumentProcessingError: If extraction fails
        """
        try:
            logger.info("Extracting text from PDF using opendataloader-pdf")

            # Use opendataloader-pdf for better table extraction
            try:
                from opendataloader_pdf import load_pdf
                
                # opendataloader-pdf processes bytes directly
                result = load_pdf(file_content)
                
                if result:
                    logger.info(f"opendataloader-pdf extracted {len(result)} characters from PDF")
                    return result
                else:
                    logger.warning("opendataloader-pdf returned no text, falling back to PyPDF2")
                        
            except ImportError:
                logger.info("opendataloader-pdf not available, using PyPDF2 fallback")
            except Exception as e:
                logger.warning(f"opendataloader-pdf extraction failed: {e}, falling back to PyPDF2")

            # Fallback to PyPDF2
            logger.info("Using PyPDF2 for PDF extraction")
            pdf_file = io.BytesIO(file_content)
            reader = PdfReader(pdf_file)

            # Extract text from all pages
            text_parts = []
            for page_num, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
                    logger.debug(f"Extracted text from page {page_num + 1}")
                except Exception as e:
                    logger.warning(
                        f"Failed to extract text from page {page_num + 1}: {str(e)}"
                    )

            if not text_parts:
                raise DocumentProcessingError("No text could be extracted from PDF")

            full_text = "\n\n".join(text_parts)
            logger.info(f"PyPDF2 extracted {len(full_text)} characters from PDF")

            return full_text

        except DocumentProcessingError:
            raise
        except Exception as e:
            error_msg = f"Failed to extract text from PDF: {str(e)}"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg) from e

    def extract_text_from_txt(self, file_content: bytes) -> str:
        """
        Extract text from TXT file.

        Args:
            file_content: TXT file content as bytes

        Returns:
            str: Extracted text

        Raises:
            DocumentProcessingError: If extraction fails
        """
        try:
            logger.info("Extracting text from TXT")

            # Try different encodings
            encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]

            for encoding in encodings:
                try:
                    text = file_content.decode(encoding)
                    logger.info(f"Successfully decoded TXT with {encoding} encoding")
                    logger.info(f"Extracted {len(text)} characters from TXT")
                    return text
                except UnicodeDecodeError:
                    continue

            # If all encodings fail
            raise DocumentProcessingError(
                f"Failed to decode TXT file with any of: {', '.join(encodings)}"
            )

        except DocumentProcessingError:
            raise
        except Exception as e:
            error_msg = f"Failed to extract text from TXT: {str(e)}"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg) from e

    def extract_text_from_docx(self, file_content: bytes) -> str:
        """
        Extract text from DOCX file including tables.

        Args:
            file_content: DOCX file content as bytes

        Returns:
            str: Extracted text with tables formatted as text

        Raises:
            DocumentProcessingError: If extraction fails
        """
        try:
            logger.info("Extracting text from DOCX")

            # Create file-like object from bytes
            docx_file = io.BytesIO(file_content)

            # Read DOCX
            doc = DocxDocument(docx_file)

            # Extract text from all paragraphs
            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Extract tables
            if doc.tables:
                logger.info(f"Found {len(doc.tables)} tables in DOCX")
                for table_idx, table in enumerate(doc.tables, 1):
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)

                    if table_data:
                        table_text = self._format_table_as_text(table_data)
                        text_parts.append(f"\n[Table {table_idx}]\n{table_text}\n")

            if not text_parts:
                raise DocumentProcessingError("No text could be extracted from DOCX")

            full_text = "\n\n".join(text_parts)
            logger.info(f"Successfully extracted {len(full_text)} characters from DOCX")

            return full_text

        except DocumentProcessingError:
            raise
        except Exception as e:
            error_msg = f"Failed to extract text from DOCX: {str(e)}"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg) from e

    def extract_text_from_hwp(self, file_content: bytes) -> str:
        """
        Extract text from HWP file (Korean Hangul Word Processor format).

        HWP files are OLE2 compound documents. This implementation uses olefile
        to parse the structure and extract text from the BodyText streams.

        Args:
            file_content: HWP file content as bytes

        Returns:
            str: Extracted text

        Raises:
            DocumentProcessingError: If extraction fails or file is corrupted/password-protected
        """
        try:
            logger.info("Extracting text from HWP")

            # Create file-like object from bytes
            hwp_file = io.BytesIO(file_content)

            # Check if it's a valid OLE file
            if not olefile.isOleFile(hwp_file):
                raise DocumentProcessingError(
                    "Invalid HWP file: not a valid OLE2 compound document"
                )

            # Open OLE file
            hwp_file.seek(0)  # Reset position
            ole = olefile.OleFileIO(hwp_file)

            try:
                # Check for password protection
                # HWP files store encryption info in DocInfo stream
                if ole.exists("DocInfo"):
                    doc_info = ole.openstream("DocInfo").read()
                    # Simple check for encryption flag (this is a heuristic)
                    if b"Password" in doc_info or b"Encrypt" in doc_info:
                        raise DocumentProcessingError(
                            "HWP file appears to be password-protected. "
                            "Please remove password protection and try again."
                        )

                # Extract text from BodyText sections
                text_parts = []
                table_count = 0

                # HWP files store text in BodyText/Section* streams
                for entry in ole.listdir():
                    entry_path = "/".join(entry)

                    # Look for BodyText sections
                    if entry_path.startswith("BodyText/Section"):
                        try:
                            stream = ole.openstream(entry)
                            content = stream.read()

                            # HWP text is typically UTF-16LE encoded
                            # Try to decode and extract readable text
                            try:
                                # Try UTF-16LE first (most common)
                                decoded = content.decode("utf-16le", errors="ignore")
                            except:
                                # Fallback to UTF-8
                                decoded = content.decode("utf-8", errors="ignore")

                            # Clean up the text (remove null bytes and control characters)
                            cleaned = "".join(
                                char
                                for char in decoded
                                if char.isprintable() or char in "\n\r\t "
                            )

                            # Try to detect and format table-like structures
                            # HWP tables often have tab-separated values or specific markers
                            if "\t" in cleaned:
                                # Potential table detected
                                lines = cleaned.split("\n")
                                table_lines = []
                                current_table = []

                                for line in lines:
                                    if "\t" in line:
                                        # This line might be part of a table
                                        cells = [
                                            cell.strip()
                                            for cell in line.split("\t")
                                            if cell.strip()
                                        ]
                                        if cells:
                                            current_table.append(cells)
                                    else:
                                        # Not a table line
                                        if current_table and len(current_table) > 1:
                                            # We have accumulated table rows
                                            table_count += 1
                                            table_text = self._format_table_as_text(
                                                current_table
                                            )
                                            table_lines.append(
                                                f"\n[Table {table_count}]\n{table_text}\n"
                                            )
                                            current_table = []
                                        if line.strip():
                                            table_lines.append(line.strip())

                                # Check for remaining table
                                if current_table and len(current_table) > 1:
                                    table_count += 1
                                    table_text = self._format_table_as_text(
                                        current_table
                                    )
                                    table_lines.append(
                                        f"\n[Table {table_count}]\n{table_text}\n"
                                    )

                                if table_lines:
                                    text_parts.append("\n".join(table_lines))
                            elif cleaned.strip():
                                text_parts.append(cleaned.strip())

                            logger.debug(f"Extracted text from {entry_path}")

                        except Exception as e:
                            logger.warning(
                                f"Failed to extract from {entry_path}: {str(e)}"
                            )
                            continue

                if not text_parts:
                    raise DocumentProcessingError(
                        "No text could be extracted from HWP file. "
                        "The file may be corrupted or use an unsupported HWP version."
                    )

                full_text = "\n\n".join(text_parts)
                logger.info(
                    f"Successfully extracted {len(full_text)} characters from HWP"
                )

                return full_text

            finally:
                ole.close()

        except DocumentProcessingError:
            raise
        except Exception as e:
            error_msg = f"Failed to extract text from HWP: {str(e)}"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg) from e

    def extract_text_from_hwpx(self, file_content: bytes) -> str:
        """
        Extract text from HWPX file (Korean Hangul XML format).

        HWPX is a ZIP-based XML format similar to DOCX, introduced in Hancom Office 2014+.
        Text content is stored in Contents/section*.xml files.

        Args:
            file_content: HWPX file content as bytes

        Returns:
            str: Extracted text

        Raises:
            DocumentProcessingError: If extraction fails or file is corrupted/password-protected
        """
        try:
            logger.info("Extracting text from HWPX")

            # Create file-like object from bytes
            hwpx_file = io.BytesIO(file_content)

            # Check if it's a valid ZIP file
            if not zipfile.is_zipfile(hwpx_file):
                raise DocumentProcessingError(
                    "Invalid HWPX file: not a valid ZIP archive"
                )

            # Open ZIP file
            hwpx_file.seek(0)  # Reset position

            with zipfile.ZipFile(hwpx_file, "r") as zip_ref:
                # Check for password protection
                for file_info in zip_ref.filelist:
                    if file_info.flag_bits & 0x1:  # Bit 0 indicates encryption
                        raise DocumentProcessingError(
                            "HWPX file is password-protected. "
                            "Please remove password protection and try again."
                        )

                # List all files in the archive
                file_list = zip_ref.namelist()

                # Find section XML files in Contents directory
                section_files = sorted(
                    [
                        f
                        for f in file_list
                        if f.startswith("Contents/section") and f.endswith(".xml")
                    ]
                )

                if not section_files:
                    raise DocumentProcessingError(
                        "No section files found in HWPX archive. "
                        "The file may be corrupted."
                    )

                text_parts = []
                table_count = 0

                # Extract text from each section
                for section_file in section_files:
                    try:
                        # Read XML content
                        xml_content = zip_ref.read(section_file)

                        # Parse XML
                        try:
                            root = ET.fromstring(xml_content)
                        except ET.ParseError as e:
                            logger.warning(f"Failed to parse {section_file}: {str(e)}")
                            continue

                        # Extract tables first (HWPX uses hp:tbl tag for tables)
                        # Define namespace if needed
                        namespaces = {
                            "hp": "http://www.hancom.co.kr/hwpml/2011/paragraph"
                        }

                        # Try to find tables with and without namespace
                        tables = root.findall(".//hp:tbl", namespaces) or root.findall(
                            ".//tbl"
                        )

                        if tables:
                            logger.info(f"Found {len(tables)} tables in {section_file}")
                            for table_elem in tables:
                                table_count += 1
                                table_data = self._extract_hwpx_table(
                                    table_elem, namespaces
                                )
                                if table_data:
                                    table_text = self._format_table_as_text(table_data)
                                    text_parts.append(
                                        f"\n[Table {table_count}]\n{table_text}\n"
                                    )

                        # Extract regular text from all text nodes
                        # HWPX uses various tags for text content, primarily <hp:t> tags
                        section_text = self._extract_text_from_xml(root)

                        if section_text.strip():
                            text_parts.append(section_text.strip())
                            logger.debug(f"Extracted text from {section_file}")

                    except Exception as e:
                        logger.warning(f"Failed to process {section_file}: {str(e)}")
                        continue

                if not text_parts:
                    raise DocumentProcessingError(
                        "No text could be extracted from HWPX file. "
                        "The file may be corrupted or empty."
                    )

                full_text = "\n\n".join(text_parts)
                logger.info(
                    f"Successfully extracted {len(full_text)} characters from HWPX"
                )

                return full_text

        except DocumentProcessingError:
            raise
        except Exception as e:
            error_msg = f"Failed to extract text from HWPX: {str(e)}"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg) from e

    def _extract_hwpx_table(self, table_elem: ET.Element, namespaces: dict) -> list:
        """
        Extract table data from HWPX table element.

        Args:
            table_elem: XML table element
            namespaces: XML namespaces dictionary

        Returns:
            list: Table data as list of lists
        """
        table_data = []

        try:
            # HWPX tables have tr (table row) and tc (table cell) elements
            rows = table_elem.findall(".//hp:tr", namespaces) or table_elem.findall(
                ".//tr"
            )

            for row in rows:
                row_data = []
                cells = row.findall(".//hp:tc", namespaces) or row.findall(".//tc")

                for cell in cells:
                    # Extract text from cell
                    cell_text = self._extract_text_from_xml(cell)
                    row_data.append(cell_text.strip() if cell_text else "")

                if row_data:
                    table_data.append(row_data)

        except Exception as e:
            logger.warning(f"Failed to extract HWPX table: {e}")

        return table_data

    def _extract_text_from_xml(self, element: ET.Element) -> str:
        """
        Recursively extract text from XML element and its children.

        Args:
            element: XML element to extract text from

        Returns:
            str: Extracted text
        """
        text_parts = []

        # Get text from current element
        if element.text:
            text_parts.append(element.text)

        # Recursively process children
        for child in element:
            child_text = self._extract_text_from_xml(child)
            if child_text:
                text_parts.append(child_text)

            # Get tail text (text after the child element)
            if child.tail:
                text_parts.append(child.tail)

        return " ".join(text_parts)


    
    def _format_table_as_text(self, table: list) -> str:
        """
        Format a table (list of lists) as readable and searchable text.
        
        Enhanced for better search by repeating key information.

        Args:
            table: Table data as list of lists

        Returns:
            str: Formatted table text
        """
        if not table:
            return ""

        lines = []
        lines.append("=" * 70)
        lines.append("[TABLE START]")
        lines.append("=" * 70)

        # Check if first row is header
        if len(table) > 1:
            headers = [str(cell).strip() if cell else "" for cell in table[0]]
            
            # Process data rows
            for row_idx, row in enumerate(table[1:], 1):
                lines.append(f"\n--- Row {row_idx} ---")
                
                for header, cell in zip(headers, row):
                    if header and cell:
                        header_text = str(header).strip()
                        cell_text = str(cell).strip()
                        
                        if header_text and cell_text:
                            # Key-value format
                            lines.append(f"{header_text}: {cell_text}")
                            
                            # Repeat value for search optimization
                            lines.append(f"  → {cell_text}")
        else:
            # No header
            for row_idx, row in enumerate(table, 1):
                row_text = " | ".join(str(cell).strip() for cell in row if cell)
                if row_text:
                    lines.append(f"Row {row_idx}: {row_text}")

        lines.append("=" * 70)
        lines.append("[TABLE END]")
        lines.append("=" * 70)

        return "\n".join(lines)

    def extract_text_from_md(self, file_content: bytes) -> str:
        """
        Extract text from Markdown file.

        Args:
            file_content: MD file content as bytes

        Returns:
            str: Extracted text
        """
        try:
            # Markdown is plain text, just decode
            text = file_content.decode("utf-8", errors="ignore")
            return text.strip()
        except Exception as e:
            error_msg = f"Failed to extract text from MD: {str(e)}"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg) from e

    def extract_text_from_pptx(self, file_content: bytes) -> str:
        """
        Extract text from PPTX file including tables.

        Args:
            file_content: PPTX file content as bytes

        Returns:
            str: Extracted text from all slides with tables
        """
        try:
            from pptx import Presentation
            from io import BytesIO

            # Load presentation from bytes
            prs = Presentation(BytesIO(file_content))

            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []

                for shape in slide.shapes:
                    # Extract text from text boxes
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)

                    # Extract tables
                    if shape.has_table:
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_data.append(row_data)

                        if table_data:
                            table_text = self._format_table_as_text(table_data)
                            slide_text.append(f"\n[Table]\n{table_text}\n")

                if slide_text:
                    text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))

            return "\n\n".join(text_parts).strip()
        except ImportError:
            error_msg = "python-pptx library not installed. Install with: pip install python-pptx"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg)
        except Exception as e:
            error_msg = f"Failed to extract text from PPTX: {str(e)}"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg) from e

    def extract_text_from_ppt(self, file_content: bytes) -> str:
        """
        Extract text from PPT file (legacy format).
        Note: PPT extraction is limited. Consider converting to PPTX for better results.

        Args:
            file_content: PPT file content as bytes

        Returns:
            str: Extracted text
        """
        # For legacy PPT files, we'll try to use python-pptx which may work for some files
        # or return a message suggesting conversion
        try:
            return self.extract_text_from_pptx(file_content)
        except Exception as e:
            logger.warning(f"Failed to extract from PPT: {e}")
            return "[PPT file detected. Please convert to PPTX format for better text extraction.]"

    def extract_text(self, file_content: bytes, file_type: str) -> str:
        """
        Extract text from file based on type.

        Args:
            file_content: File content as bytes
            file_type: Type of file (pdf, txt, docx, hwp, hwpx, md, ppt, pptx)

        Returns:
            str: Extracted text

        Raises:
            ValueError: If file_type is invalid
            DocumentProcessingError: If extraction fails
        """
        if not file_content:
            raise ValueError("file_content cannot be empty")

        if file_type not in self.SUPPORTED_TYPES:
            raise ValueError(f"Unsupported file type: {file_type}")

        # For HWP/HWPX files, try PDF conversion first for better table/chart extraction
        if file_type in ["hwp", "hwpx"]:
            try:
                from backend.services.hwp_converter import get_hwp_converter
                
                converter = get_hwp_converter()
                pdf_content, error = converter.convert_to_pdf(file_content, f"temp.{file_type}")
                
                if pdf_content:
                    logger.info(f"✅ HWP→PDF conversion successful, using PDF extraction")
                    # Use PDF extraction for better table/chart recognition
                    text = self.extract_text_from_pdf(pdf_content)
                    text = self._normalize_text(text)
                    return text
                else:
                    logger.warning(f"HWP→PDF conversion failed: {error}, using native extraction")
                    
            except Exception as e:
                logger.warning(f"HWP→PDF conversion error: {e}, using native extraction")

        # Route to appropriate extraction method
        extractors = {
            "pdf": self.extract_text_from_pdf,
            "txt": self.extract_text_from_txt,
            "docx": self.extract_text_from_docx,
            "hwp": self.extract_text_from_hwp,
            "hwpx": self.extract_text_from_hwpx,
            "md": self.extract_text_from_md,
            "ppt": self.extract_text_from_ppt,
            "pptx": self.extract_text_from_pptx,
            "xlsx": self.extract_text_from_xlsx,
            "xls": self.extract_text_from_xls,
            "csv": self.extract_text_from_csv,
            "json": self.extract_text_from_json,
            # Image formats (OCR)
            "png": self.extract_text_from_image,
            "jpg": self.extract_text_from_image,
            "jpeg": self.extract_text_from_image,
            "gif": self.extract_text_from_image,
            "bmp": self.extract_text_from_image,
            "webp": self.extract_text_from_image,
        }

        extractor = extractors.get(file_type)
        if not extractor:
            raise ValueError(f"No extractor found for file type: {file_type}")
        
        # Extract text
        text = extractor(file_content)
        
        # Normalize text (remove intra-word spaces, fix formatting)
        text = self._normalize_text(text)
        
        return text

    def _normalize_text(self, text: str) -> str:
        """
        Normalize extracted text by removing intra-word spaces and fixing formatting.
        
        This is especially important for Korean documents where OCR or PDF extraction
        may introduce unwanted spaces within words.
        
        Args:
            text: Raw extracted text
            
        Returns:
            Normalized text
        """
        try:
            from backend.services.korean_text_processor import get_korean_text_processor
            
            # Get Korean text processor
            processor = get_korean_text_processor(use_morpheme_analysis=False)
            
            # Process text (includes spacing normalization)
            result = processor.process_korean_text(text)
            
            return result['normalized']
            
        except Exception as e:
            logger.warning(f"Text normalization failed: {e}. Returning original text.")
            return text
    
    def chunk_text(self, text: str, document_id: str) -> List[TextChunk]:
        """
        Split text into chunks using semantic chunking strategy.

        Args:
            text: Text to chunk
            document_id: ID of the parent document

        Returns:
            List[TextChunk]: List of text chunks

        Raises:
            ValueError: If inputs are invalid
        """
        if not text:
            raise ValueError("text cannot be empty")
        if not document_id:
            raise ValueError("document_id cannot be empty")

        # Clean text
        text = text.strip()

        if not text:
            raise ValueError("text cannot be only whitespace")

        logger.info(
            f"Chunking text of length {len(text)} with "
            f"strategy={self.chunking_strategy}, "
            f"chunk_size={self.chunk_size}, overlap={self.chunk_overlap}"
        )

        # Use semantic chunker
        try:
            chunk_texts = self.semantic_chunker.chunk_text(text)
        except Exception as e:
            logger.error(f"Semantic chunking failed: {e}, using fallback")
            # Fallback to simple chunking
            chunk_texts = self._fallback_chunking(text)

        # Convert to TextChunk objects
        chunks = []
        current_pos = 0

        for chunk_index, chunk_text in enumerate(chunk_texts):
            # Find position in original text
            start_pos = text.find(chunk_text, current_pos)
            if start_pos == -1:
                # If exact match not found, use approximate position
                start_pos = current_pos

            end_pos = start_pos + len(chunk_text)

            chunk = TextChunk(
                chunk_id=f"{document_id}_chunk_{chunk_index}",
                document_id=document_id,
                text=chunk_text,
                chunk_index=chunk_index,
                start_char=start_pos,
                end_char=end_pos,
                metadata={
                    "chunking_strategy": self.chunking_strategy,
                    "chunk_size": len(chunk_text),
                },
            )
            chunks.append(chunk)

            current_pos = end_pos

        logger.info(
            f"Created {len(chunks)} chunks from text "
            f"(avg size: {sum(len(c.text) for c in chunks) / len(chunks):.0f} chars)"
        )

        return chunks

    def _fallback_chunking(self, text: str) -> List[str]:
        """
        Fallback chunking method (simple fixed-size with sentence boundaries).

        Args:
            text: Text to chunk

        Returns:
            List of chunk texts
        """
        chunks = []
        start = 0
        text_len = len(text)

        while start < text_len:
            end = start + self.chunk_size

            if end >= text_len:
                chunks.append(text[start:].strip())
                break

            # Try to break at sentence boundary
            sentence_end = max(
                text.rfind(".", start, end),
                text.rfind("!", start, end),
                text.rfind("?", start, end),
            )

            if sentence_end > start:
                end = sentence_end + 1
            else:
                # Try word boundary
                space_pos = text.rfind(" ", start, end)
                if space_pos > start:
                    end = space_pos

            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append(chunk_text)

            start = end - self.chunk_overlap
            if start <= 0:
                start = end

        return chunks

    def extract_metadata(
        self, filename: str, file_size: int, file_type: str, text: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Extract metadata from document.

        Args:
            filename: Original filename
            file_size: File size in bytes
            file_type: Type of file
            text: Extracted text (optional, for additional metadata)

        Returns:
            Dict[str, Any]: Metadata dictionary
        """
        metadata = {
            "filename": filename,
            "file_size": file_size,
            "file_type": file_type,
            "upload_date": datetime.now().isoformat(),
            "processor_version": "1.0",
        }

        # Add text-based metadata if text is provided
        if text:
            metadata["character_count"] = len(text)
            metadata["word_count"] = len(text.split())
            metadata["line_count"] = text.count("\n") + 1

        return metadata

    async def process_document(
        self, file_content: bytes, filename: str, file_size: int
    ) -> tuple[Document, List[TextChunk]]:
        """
        Process a document end-to-end: validate, extract, chunk, extract metadata.

        Args:
            file_content: File content as bytes
            filename: Original filename
            file_size: File size in bytes

        Returns:
            tuple: (Document, List[TextChunk])

        Raises:
            ValueError: If inputs are invalid
            DocumentProcessingError: If processing fails
        """
        try:
            logger.info(f"Processing document: {filename} ({file_size} bytes)")

            # Validate file size
            self.validate_file_size(file_size)

            # Detect file type
            file_type = self.detect_file_type(filename)

            # Generate document ID
            document_id = str(uuid.uuid4())

            # Create document record
            document = Document(
                document_id=document_id,
                filename=filename,
                file_type=file_type,
                file_size=file_size,
                upload_timestamp=datetime.now(),
                processing_status="processing",
                metadata={},
            )

            # Extract text
            try:
                text = self.extract_text(file_content, file_type)
            except DocumentProcessingError as e:
                document.processing_status = "failed"
                document.error_message = str(e)
                raise

            # Chunk text
            chunks = self.chunk_text(text, document_id)

            # Extract rich metadata
            from backend.services.metadata_extractor import get_metadata_extractor

            metadata_extractor = get_metadata_extractor()
            rich_metadata = metadata_extractor.extract_metadata(
                file_content=file_content, file_type=file_type, text=text
            )

            # Store metadata
            metadata = self.extract_metadata(filename, file_size, file_type, text)
            metadata.update(rich_metadata)
            document.metadata = metadata

            # Update document status
            document.processing_status = "completed"
            document.chunk_count = len(chunks)

            logger.info(
                f"Successfully processed document {document_id}: "
                f"{len(chunks)} chunks created, metadata extracted"
            )

            return document, chunks

        except ValueError:
            raise
        except DocumentProcessingError:
            raise
        except Exception as e:
            error_msg = f"Unexpected error processing document: {str(e)}"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg) from e

    def extract_text_from_xlsx(self, file_content: bytes) -> str:
        """
        Extract text from XLSX file (Excel 2007+).

        Args:
            file_content: XLSX file content as bytes

        Returns:
            str: Extracted text with sheet names and data

        Raises:
            DocumentProcessingError: If extraction fails
        """
        try:
            logger.info("Extracting text from XLSX")
            import openpyxl
            import io

            # Create file-like object from bytes
            xlsx_file = io.BytesIO(file_content)

            # Load workbook
            workbook = openpyxl.load_workbook(xlsx_file, data_only=True)

            text_parts = []

            # Extract from each sheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                text_parts.append(f"\n[Sheet: {sheet_name}]\n")

                # Get all rows
                rows_data = []
                for row in sheet.iter_rows(values_only=True):
                    # Filter out empty rows
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(val.strip() for val in row_values):
                        rows_data.append(row_values)

                if rows_data:
                    # Format as table
                    table_text = self._format_table_as_text(rows_data)
                    text_parts.append(table_text)

            full_text = "\n\n".join(text_parts)
            logger.info(f"Successfully extracted {len(full_text)} characters from XLSX")

            return full_text

        except ImportError:
            error_msg = "openpyxl library not installed. Install with: pip install openpyxl"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg)
        except Exception as e:
            error_msg = f"Failed to extract text from XLSX: {str(e)}"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg) from e

    def extract_text_from_xls(self, file_content: bytes) -> str:
        """
        Extract text from XLS file (Excel 97-2003).

        Args:
            file_content: XLS file content as bytes

        Returns:
            str: Extracted text with sheet names and data

        Raises:
            DocumentProcessingError: If extraction fails
        """
        try:
            logger.info("Extracting text from XLS")
            import xlrd
            import io

            # Create file-like object from bytes
            xls_file = io.BytesIO(file_content)

            # Load workbook
            workbook = xlrd.open_workbook(file_contents=file_content)

            text_parts = []

            # Extract from each sheet
            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                
                text_parts.append(f"\n[Sheet: {sheet.name}]\n")

                # Get all rows
                rows_data = []
                for row_idx in range(sheet.nrows):
                    row_values = []
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        row_values.append(str(cell.value) if cell.value else "")
                    
                    if any(val.strip() for val in row_values):
                        rows_data.append(row_values)

                if rows_data:
                    # Format as table
                    table_text = self._format_table_as_text(rows_data)
                    text_parts.append(table_text)

            full_text = "\n\n".join(text_parts)
            logger.info(f"Successfully extracted {len(full_text)} characters from XLS")

            return full_text

        except ImportError:
            error_msg = "xlrd library not installed. Install with: pip install xlrd"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg)
        except Exception as e:
            error_msg = f"Failed to extract text from XLS: {str(e)}"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg) from e

    def extract_text_from_csv(self, file_content: bytes) -> str:
        """
        Extract text from CSV file.

        Args:
            file_content: CSV file content as bytes

        Returns:
            str: Extracted text formatted as table

        Raises:
            DocumentProcessingError: If extraction fails
        """
        try:
            logger.info("Extracting text from CSV")
            import csv
            import io

            # Try different encodings
            encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252", "euc-kr"]
            
            text = None
            for encoding in encodings:
                try:
                    text = file_content.decode(encoding)
                    logger.info(f"Successfully decoded CSV with {encoding} encoding")
                    break
                except UnicodeDecodeError:
                    continue

            if text is None:
                raise DocumentProcessingError(
                    f"Failed to decode CSV file with any of: {', '.join(encodings)}"
                )

            # Parse CSV
            csv_file = io.StringIO(text)
            reader = csv.reader(csv_file)

            rows_data = []
            for row in reader:
                if any(cell.strip() for cell in row):
                    rows_data.append(row)

            if not rows_data:
                raise DocumentProcessingError("No data found in CSV file")

            # Format as table
            table_text = self._format_table_as_text(rows_data)
            
            full_text = f"[CSV Data]\n{table_text}"
            logger.info(f"Successfully extracted {len(full_text)} characters from CSV")

            return full_text

        except DocumentProcessingError:
            raise
        except Exception as e:
            error_msg = f"Failed to extract text from CSV: {str(e)}"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg) from e

    def extract_text_from_json(self, file_content: bytes) -> str:
        """
        Extract text from JSON file.

        Args:
            file_content: JSON file content as bytes

        Returns:
            str: Extracted text formatted as readable structure

        Raises:
            DocumentProcessingError: If extraction fails
        """
        try:
            logger.info("Extracting text from JSON")
            import json

            # Try different encodings
            encodings = ["utf-8", "utf-8-sig", "latin-1"]
            
            text = None
            for encoding in encodings:
                try:
                    text = file_content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue

            if text is None:
                raise DocumentProcessingError(
                    f"Failed to decode JSON file with any of: {', '.join(encodings)}"
                )

            # Parse JSON
            data = json.loads(text)

            # Convert to readable text
            full_text = self._format_json_as_text(data)
            
            logger.info(f"Successfully extracted {len(full_text)} characters from JSON")

            return full_text

        except json.JSONDecodeError as e:
            error_msg = f"Invalid JSON format: {str(e)}"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg) from e
        except DocumentProcessingError:
            raise
        except Exception as e:
            error_msg = f"Failed to extract text from JSON: {str(e)}"
            logger.error(error_msg)
            raise DocumentProcessingError(error_msg) from e

    def _format_json_as_text(self, data, indent: int = 0) -> str:
        """
        Format JSON data as readable text.

        Args:
            data: JSON data (dict, list, or primitive)
            indent: Current indentation level

        Returns:
            str: Formatted text
        """
        lines = []
        indent_str = "  " * indent

        if isinstance(data, dict):
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    lines.append(f"{indent_str}{key}:")
                    lines.append(self._format_json_as_text(value, indent + 1))
                else:
                    lines.append(f"{indent_str}{key}: {value}")
        
        elif isinstance(data, list):
            for idx, item in enumerate(data):
                if isinstance(item, (dict, list)):
                    lines.append(f"{indent_str}[{idx}]:")
                    lines.append(self._format_json_as_text(item, indent + 1))
                else:
                    lines.append(f"{indent_str}- {item}")
        
        else:
            lines.append(f"{indent_str}{data}")

        return "\n".join(lines)

    def extract_text_from_image(self, file_content: bytes) -> str:
        """
        Extract text from image file using PaddleOCR.

        Supports: PNG, JPG, JPEG, GIF, BMP, WEBP

        Uses PaddleOCR (Official GitHub):
        - Text detection and recognition
        - Text direction classification
        - Multi-language support (Korean optimized)
        - GPU acceleration
        - 95%+ accuracy for printed text

        Args:
            file_content: Image file content as bytes

        Returns:
            str: Extracted text with metadata

        Raises:
            DocumentProcessingError: If extraction fails
        """
        try:
            # Use PaddleOCR Advanced for full document understanding
            try:
                from backend.services.paddleocr_advanced import get_paddleocr_advanced
                from backend.config import settings
                
                logger.info("Extracting text from image using PaddleOCR Advanced")
                
                # Get PaddleOCR Advanced processor
                processor = get_paddleocr_advanced(
                    lang=settings.PADDLEOCR_LANG,
                    use_gpu=settings.PADDLEOCR_USE_GPU,
                    enable_table_recognition=True,
                    enable_layout_analysis=True
                )
                
                if not processor or not processor.ocr_available:
                    logger.warning("PaddleOCR not available")
                    return "[Image content - PaddleOCR not installed. Install with: pip install paddlepaddle-gpu paddleocr[all]]"
                
                # Process document with full capabilities
                result = processor.process_document(
                    file_content,
                    extract_tables=True,
                    analyze_layout=True,
                    cls=True
                )
                
                # Build comprehensive text output
                text_parts = []
                
                # 1. Main text
                if result['text']:
                    text_parts.append("=" * 70)
                    text_parts.append("[EXTRACTED TEXT - PP-OCRv5]")
                    text_parts.append("=" * 70)
                    text_parts.append(result['text'])
                    text_parts.append("")
                
                # 2. Tables (PP-StructureV3)
                if result['tables']:
                    text_parts.append("=" * 70)
                    text_parts.append(f"[TABLES - PP-StructureV3] ({len(result['tables'])} detected)")
                    text_parts.append("=" * 70)
                    
                    for i, table in enumerate(result['tables'], 1):
                        text_parts.append(f"\n--- Table {i} ---")
                        text_parts.append(f"Confidence: {table['confidence']:.2%}")
                        
                        # Add HTML table structure if available
                        if table.get('html'):
                            text_parts.append("HTML Structure:")
                            text_parts.append(table['html'][:500])  # First 500 chars
                        
                        # Add cell information
                        if table.get('cells'):
                            text_parts.append(f"Cells: {len(table['cells'])}")
                        
                        text_parts.append("")
                
                # 3. Layout Analysis
                if result['layout']:
                    text_parts.append("=" * 70)
                    text_parts.append(f"[LAYOUT ANALYSIS] ({len(result['layout'])} regions)")
                    text_parts.append("=" * 70)
                    
                    # Group by type
                    layout_summary = {}
                    for region in result['layout']:
                        region_type = region['type']
                        layout_summary[region_type] = layout_summary.get(region_type, 0) + 1
                    
                    for region_type, count in sorted(layout_summary.items()):
                        text_parts.append(f"- {region_type}: {count} region(s)")
                    
                    text_parts.append("")
                
                # 4. Statistics
                stats = result['stats']
                text_parts.append("=" * 70)
                text_parts.append("[OCR STATISTICS]")
                text_parts.append("=" * 70)
                text_parts.append(f"Text Boxes: {stats['num_text_boxes']}")
                text_parts.append(f"Tables: {stats['num_tables']}")
                text_parts.append(f"Layout Regions: {stats['num_layout_regions']}")
                text_parts.append(f"Total Characters: {stats['total_characters']}")
                text_parts.append(f"Average Confidence: {stats['avg_confidence']:.2%}")
                text_parts.append("=" * 70)
                
                final_text = "\n".join(text_parts)
                
                logger.info(
                    f"✅ PaddleOCR Advanced: {stats['total_characters']} chars, "
                    f"{stats['num_tables']} tables, {stats['num_layout_regions']} regions "
                    f"(confidence: {stats['avg_confidence']:.2%})"
                )
                
                return final_text
                
            except ImportError as e:
                logger.warning(f"PaddleOCR import failed: {e}")
                logger.warning(
                    "Install with: "
                    "pip install paddlepaddle-gpu==3.0.0b1 -i https://www.paddlepaddle.org.cn/packages/stable/cu118/ && "
                    "pip install paddleocr>=2.7.0"
                )
            except Exception as e:
                logger.error(f"PaddleOCR processing failed: {e}")
                import traceback
                traceback.print_exc()
            
            # Fallback: Return placeholder
            logger.warning("PaddleOCR not available, returning placeholder")
            return "[Image content - PaddleOCR not installed. Install with: pip install paddlepaddle-gpu==3.0.0b1 paddleocr>=2.7.0]"

        except Exception as e:
            error_msg = f"Failed to process image: {str(e)}"
            logger.error(error_msg)
            # Don't raise error, return placeholder instead
            return f"[Image content - processing failed: {str(e)}]"

    def _format_table_as_text(self, rows: List[List[str]]) -> str:
        """
        Format table data as readable text.

        Args:
            rows: List of rows, each row is a list of cell values

        Returns:
            str: Formatted table text
        """
        if not rows:
            return ""

        # Calculate column widths
        col_widths = []
        for col_idx in range(len(rows[0])):
            max_width = max(
                len(str(row[col_idx])) if col_idx < len(row) else 0
                for row in rows
            )
            col_widths.append(min(max_width, 50))  # Cap at 50 chars

        # Format rows
        formatted_rows = []
        for row in rows:
            formatted_cells = []
            for col_idx, cell in enumerate(row):
                if col_idx < len(col_widths):
                    cell_str = str(cell)[:col_widths[col_idx]]
                    formatted_cells.append(cell_str.ljust(col_widths[col_idx]))
            formatted_rows.append(" | ".join(formatted_cells))

        return "\n".join(formatted_rows)

    def _format_json_as_text(self, data: Any, indent: int = 0) -> str:
        """
        Format JSON data as readable text.

        Args:
            data: JSON data (dict, list, or primitive)
            indent: Current indentation level

        Returns:
            str: Formatted text
        """
        import json
        
        # For simple types, use JSON dump
        if isinstance(data, (str, int, float, bool)) or data is None:
            return json.dumps(data, ensure_ascii=False)

        # For complex types, format nicely
        if isinstance(data, dict):
            lines = []
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    lines.append(f"{'  ' * indent}{key}:")
                    lines.append(self._format_json_as_text(value, indent + 1))
                else:
                    lines.append(
                        f"{'  ' * indent}{key}: {json.dumps(value, ensure_ascii=False)}"
                    )
            return "\n".join(lines)

        elif isinstance(data, list):
            if not data:
                return "[]"
            
            lines = []
            for idx, item in enumerate(data):
                if isinstance(item, (dict, list)):
                    lines.append(f"{'  ' * indent}[{idx}]:")
                    lines.append(self._format_json_as_text(item, indent + 1))
                else:
                    lines.append(
                        f"{'  ' * indent}[{idx}]: {json.dumps(item, ensure_ascii=False)}"
                    )
            return "\n".join(lines)

        return json.dumps(data, ensure_ascii=False, indent=2)
